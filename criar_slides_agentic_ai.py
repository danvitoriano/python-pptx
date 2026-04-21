import argparse
import json
import re
import sys
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

# Configuracoes de estilo
COR_FUNDO = RGBColor(0, 0, 0)
COR_TITULO = RGBColor(255, 255, 255)
COR_TEXTO = RGBColor(204, 204, 204)
COR_DESTAQUE = RGBColor(102, 179, 255)
FONTE_PRINCIPAL = "Segoe UI"
FONTE_TITULO = "Segoe UI"
FONTE_SUBTITULO = "Segoe UI"

LAYOUTS_SUPORTADOS = {
    "comparativo",
    "conceitos",
    "titulo",
    "title_left_text_right",
    "title_top_bullets",
    "title_top_grid_2x2",
    "title_top_text_block",
}

THEME_DEFAULT = {
    "profile": "default",
    "slide": {
        "ratio": "16:9",
        "width_in": 13.333,
        "height_in": 7.5,
    },
    "colors": {
        "background": "#000000",
        "title": "#FFFFFF",
        "body": "#CCCCCC",
        "accent": "#66B3FF",
        "card_bg": "#1A1A1A",
        "card_border": "#646464",
    },
    "fonts": {
        "title_family": "Segoe UI",
        "subtitle_family": "Segoe UI",
        "body_family": "Segoe UI",
        "title_size_max": 40,
        "title_size_mid": 32,
        "title_size_small": 28,
        "title_size_min": 24,
        "subtitle_size": 24,
        "body_size": 14,
    },
    "layout": {
        "margin_x_in": 0.8,
        "title_top_in": 0.5,
        "subtitle_top_in": 1.3,
        "content_top_in": 2.0,
        "column_gap_in": 0.5,
        "bottom_block_top_in": 4.8,
        "bottom_block_height_in": 2.0,
        "max_content_top_in": 2.2,
        "max_left_title_width_in": 3.2,
        "max_left_content_width_in": 8.4,
    },
    "parsing": {
        "strip_page_prefix": True,
        "ignore_sections": ["Notas Estruturais para Python-PPTX"],
    },
    "title": {
        "accent_words": [
            "DESIGN VISUAL",
            "SLIDE EM MAIUSCULO",
            "SLIDE EM MAIUSCULA",
            "PRATICAS",
            "BOAS PRATICAS",
            "COMPARATIVO",
            "RESULTADO",
        ],
        "force_uppercase": False,
    },
}

THEME_ACTIVE = deepcopy(THEME_DEFAULT)


def _ajustar_text_frame(tf):
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE


def _tamanho_titulo(texto):
    fonts = THEME_ACTIVE["fonts"]
    tamanho = len(texto)
    if tamanho > 85:
        return Pt(fonts["title_size_min"])
    if tamanho > 65:
        return Pt(fonts["title_size_small"])
    if tamanho > 45:
        return Pt(fonts["title_size_mid"])
    return Pt(fonts["title_size_max"])


def _decorar_titulo(texto):
    mapa = [
        ("termostato", "🌡️"),
        ("climatiza", "🧠"),
        ("estado", "🧠"),
        ("objetivo", "🎯"),
        ("anatomia", "🔧"),
        ("sensores", "📡"),
        ("atuadores", "⚙️"),
        ("motor", "🔄"),
    ]
    base = texto.strip()
    low = base.lower()
    for termo, emoji in mapa:
        if termo in low and emoji not in base:
            return f"{emoji} {base}"
    return base


def _split_title_accent(texto):
    clean = (texto or "").strip()
    upper = clean.upper()
    accent_words = THEME_ACTIVE.get("title", {}).get("accent_words", [])
    for item in accent_words:
        up_item = str(item).upper().strip()
        pos = upper.find(up_item)
        if pos > -1:
            left = clean[:pos].strip()
            return left, clean[pos : pos + len(up_item)].strip(), clean[pos + len(up_item) :].strip()

    if ":" in clean:
        left, right = clean.split(":", 1)
        return left.strip(), right.strip(), ""

    marker = " DE "
    pos = upper.rfind(marker)
    if pos > -1:
        left = clean[:pos].strip()
        right = clean[pos + len(marker) :].strip()
        return left, f"DE {right}", ""

    return clean, "", ""


def _formatar_titulo_slide(texto):
    value = (texto or "").strip()
    if THEME_ACTIVE.get("title", {}).get("force_uppercase", False):
        return value.upper()
    return value


def _adicionar_titulo_max(slide, texto, top=None):
    texto = _formatar_titulo_slide(texto)
    mx = THEME_ACTIVE["layout"]["margin_x_in"]
    largura = THEME_ACTIVE["slide"]["width_in"] - (mx * 2)
    top = THEME_ACTIVE["layout"]["title_top_in"] if top is None else top
    textbox = slide.shapes.add_textbox(Inches(mx), Inches(top), Inches(largura), Inches(2.15))
    tf = textbox.text_frame
    tf.clear()
    _ajustar_text_frame(tf)
    tf.paragraphs[0].clear()

    principal, destaque, sufixo = _split_title_accent(texto)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = principal
    run.font.name = FONTE_TITULO
    run.font.size = _tamanho_titulo(texto)
    run.font.bold = False
    run.font.color.rgb = COR_TITULO

    if destaque:
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = destaque
        run2.font.name = FONTE_TITULO
        run2.font.size = _tamanho_titulo(texto)
        run2.font.bold = False
        run2.font.color.rgb = COR_DESTAQUE
        if sufixo:
            run3 = p2.add_run()
            run3.text = f" {sufixo}"
            run3.font.name = FONTE_TITULO
            run3.font.size = _tamanho_titulo(texto)
            run3.font.bold = False
            run3.font.color.rgb = COR_TITULO
    return textbox


def _merge_dict(base, override):
    for key, value in override.items():
        if key in base and isinstance(base[key], dict) and isinstance(value, dict):
            _merge_dict(base[key], value)
        else:
            base[key] = value
    return base


def _parse_hex_color(valor):
    valor = valor.strip().lstrip("#")
    if len(valor) != 6:
        raise ValueError(f"Cor invalida '{valor}'. Use formato #RRGGBB.")
    return RGBColor(int(valor[0:2], 16), int(valor[2:4], 16), int(valor[4:6], 16))


def _load_theme(theme_path=None, profile=None):
    tema = deepcopy(THEME_DEFAULT)
    if theme_path:
        path = Path(theme_path)
        raw = path.read_text(encoding="utf-8")
        try:
            dados = json.loads(raw)
        except json.JSONDecodeError as exc:
            raise ValueError(f"Arquivo de tema invalido ({path}): {exc}") from exc
        _merge_dict(tema, dados)
    if profile:
        tema["profile"] = profile
        if profile == "premium":
            _merge_dict(
                tema,
                {
                    "fonts": {
                        "title_size_max": 44,
                        "title_size_mid": 36,
                        "title_size_small": 30,
                        "title_size_min": 24,
                        "subtitle_size": 24,
                        "body_size": 13,
                    },
                    "layout": {
                        "column_gap_in": 0.6,
                        "content_top_in": 1.95,
                        "bottom_block_top_in": 4.75,
                    },
                },
            )
        if profile == "max":
            _merge_dict(
                tema,
                {
                    "fonts": {
                        "title_size_max": 40,
                        "title_size_mid": 35,
                        "title_size_small": 29,
                        "title_size_min": 24,
                        "subtitle_size": 20,
                        "body_size": 15,
                    },
                    "colors": {
                        "body": "#C9C9C9",
                        "accent": "#FF2E93",
                        "card_bg": "#000000",
                        "card_border": "#000000",
                    },
                    "layout": {
                        "margin_x_in": 0.85,
                        "title_top_in": 0.45,
                        "subtitle_top_in": 1.45,
                        "content_top_in": 2.0,
                        "column_gap_in": 0.45,
                        "max_content_top_in": 2.85,
                        "max_left_title_width_in": 3.3,
                        "max_left_content_width_in": 8.3,
                    },
                    "title": {
                        "accent_words": [
                            "DESIGN VISUAL",
                            "SLIDE EM MAIUSCULO",
                            "SLIDE EM MAIUSCULA",
                            "RESULTADOS",
                            "PRATICAS",
                            "REGRAS",
                        ],
                        "force_uppercase": True,
                    },
                },
            )
    return tema


def _apply_theme(tema):
    global THEME_ACTIVE, COR_FUNDO, COR_TITULO, COR_TEXTO, COR_DESTAQUE
    global FONTE_PRINCIPAL, FONTE_TITULO, FONTE_SUBTITULO

    THEME_ACTIVE = tema
    COR_FUNDO = _parse_hex_color(tema["colors"]["background"])
    COR_TITULO = _parse_hex_color(tema["colors"]["title"])
    COR_TEXTO = _parse_hex_color(tema["colors"]["body"])
    COR_DESTAQUE = _parse_hex_color(tema["colors"]["accent"])
    FONTE_TITULO = tema["fonts"]["title_family"]
    FONTE_SUBTITULO = tema["fonts"]["subtitle_family"]
    FONTE_PRINCIPAL = tema["fonts"]["body_family"]


def _should_ignore_section(title):
    ignores = THEME_ACTIVE.get("parsing", {}).get("ignore_sections", [])
    alvo = (title or "").strip().lower()
    return any(alvo == str(item).strip().lower() for item in ignores)


def configurar_slide(slide):
    """Aplica fundo preto ao slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COR_FUNDO


def adicionar_titulo(slide, texto, top=None):
    """Adiciona titulo principal."""
    if THEME_ACTIVE.get("profile") == "max":
        return _adicionar_titulo_max(slide, texto, top=top)

    texto = _formatar_titulo_slide(texto)
    mx = THEME_ACTIVE["layout"]["margin_x_in"]
    largura = THEME_ACTIVE["slide"]["width_in"] - (mx * 2)
    top = THEME_ACTIVE["layout"]["title_top_in"] if top is None else top
    textbox = slide.shapes.add_textbox(Inches(mx), Inches(top), Inches(largura), Inches(1.2))
    tf = textbox.text_frame
    tf.clear()
    _ajustar_text_frame(tf)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = texto
    run.font.name = FONTE_TITULO
    run.font.size = _tamanho_titulo(texto)
    run.font.bold = True
    run.font.color.rgb = COR_TITULO
    p.alignment = PP_ALIGN.LEFT
    return textbox


def adicionar_subtitulo(slide, texto, top=None):
    """Adiciona subtitulo."""
    mx = THEME_ACTIVE["layout"]["margin_x_in"]
    largura = THEME_ACTIVE["slide"]["width_in"] - (mx * 2)
    top = THEME_ACTIVE["layout"]["subtitle_top_in"] if top is None else top
    textbox = slide.shapes.add_textbox(Inches(mx), Inches(top), Inches(largura), Inches(0.9))
    tf = textbox.text_frame
    tf.clear()
    _ajustar_text_frame(tf)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = texto
    run.font.name = FONTE_SUBTITULO
    run.font.size = Pt(THEME_ACTIVE["fonts"]["subtitle_size"])
    run.font.color.rgb = COR_TEXTO
    p.alignment = PP_ALIGN.LEFT
    return textbox


def adicionar_lista(slide, itens, left, top, width, height, titulo=None, fonte_titulo=20, fonte_item=16):
    """Adiciona caixa com lista de itens."""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.clear()
    _ajustar_text_frame(tf)

    if titulo:
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = _decorar_titulo(titulo)
        run.font.name = FONTE_TITULO
        run.font.size = Pt(fonte_titulo)
        run.font.bold = True
        run.font.color.rgb = COR_TITULO

    for item in itens:
        p = tf.add_paragraph()
        p.space_before = Pt(8)
        run = p.add_run()
        run.text = f"• {item}"
        run.font.name = FONTE_PRINCIPAL
        run.font.size = Pt(fonte_item)
        run.font.color.rgb = COR_TEXTO

    return textbox


def adicionar_caixa_conceito(slide, titulo, descricao, left, top, width, height):
    """Adiciona caixa de conceito com titulo e descricao."""
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _parse_hex_color(THEME_ACTIVE["colors"]["card_bg"])
    shape.line.color.rgb = _parse_hex_color(THEME_ACTIVE["colors"]["card_border"])
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.clear()
    _ajustar_text_frame(tf)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = _decorar_titulo(titulo)
    run.font.name = FONTE_TITULO
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = COR_DESTAQUE

    p = tf.add_paragraph()
    p.space_before = Pt(10)
    run = p.add_run()
    run.text = descricao
    run.font.name = FONTE_PRINCIPAL
    run.font.size = Pt(THEME_ACTIVE["fonts"]["body_size"])
    run.font.color.rgb = COR_TEXTO
    p.alignment = PP_ALIGN.LEFT

    return shape


def _limpar_markdown_inline(texto):
    texto = re.sub(r"`([^`]+)`", r"\1", texto)
    texto = re.sub(r"\*\*([^*]+)\*\*", r"\1", texto)
    texto = re.sub(r"\*([^*]+)\*", r"\1", texto)
    return texto.strip()


def _remover_prefixo_pagina(texto):
    """
    Remove prefixos como 'Pagina 2:' ou 'Página 2:' para nao poluir titulos.
    """
    if not THEME_ACTIVE.get("parsing", {}).get("strip_page_prefix", True):
        return texto.strip()
    return re.sub(r"^\s*p[aá]gina\s*\d+\s*:\s*", "", texto, flags=re.IGNORECASE).strip()


def _normalizar_rotulo(texto):
    txt = texto.strip().lower()
    txt = txt.replace("í", "i").replace("ã", "a").replace("â", "a").replace("ç", "c")
    txt = re.sub(r"[^a-z]", "", txt)
    return txt


def _descricao_secao(secao):
    itens = secao.get("itens", [])
    if not itens:
        return ""
    return " ".join(itens).strip()


def _word_count(texto):
    return len(re.findall(r"\w+", texto or ""))


def _inferir_layout_generico(slide, secoes_genericas):
    if not secoes_genericas:
        slide["layout"] = "titulo"
        return

    if THEME_ACTIVE.get("profile") == "max":
        if len(secoes_genericas) >= 4:
            slide["layout"] = "title_top_grid_2x2"
            for secao in secoes_genericas[:4]:
                itens = secao.get("itens") or [secao["titulo"]]
                slide["listas"].append({"titulo": secao["titulo"], "itens": itens})
            if len(secoes_genericas) > 4:
                extra = []
                for secao in secoes_genericas[4:]:
                    extra.extend(secao.get("itens") or [secao["titulo"]])
                if extra:
                    slide["subtitle"] = " ".join(extra)[:180]
            return

        if len(secoes_genericas) == 1:
            secao = secoes_genericas[0]
            texto = " ".join(secao.get("itens", []))
            if _word_count(texto) > 45:
                slide["layout"] = "title_top_text_block"
                slide["listas"].append({"titulo": secao["titulo"], "itens": [texto]})
                return

        if len(secoes_genericas) in (2, 3):
            slide["layout"] = "title_top_bullets"
            for secao in secoes_genericas:
                itens = secao.get("itens") or [secao["titulo"]]
                slide["listas"].append({"titulo": secao["titulo"], "itens": itens})
            return

        slide["layout"] = "title_left_text_right"
        for secao in secoes_genericas:
            itens = secao.get("itens") or [secao["titulo"]]
            slide["listas"].append({"titulo": secao["titulo"], "itens": itens})
        return

    titulos = [s["titulo"].lower() for s in secoes_genericas[:2]]
    parece_conceitos = len(secoes_genericas) >= 2 and "estado" in titulos[0] and "objetivo" in titulos[1]

    if parece_conceitos:
        slide["layout"] = "conceitos"
        for secao in secoes_genericas[:2]:
            descricao = _descricao_secao(secao) or secao["titulo"]
            slide["conceitos"].append({"titulo": secao["titulo"], "descricao": descricao})
        for secao in secoes_genericas[2:]:
            itens = secao.get("itens") or [secao["titulo"]]
            slide["listas"].append({"titulo": secao["titulo"], "itens": itens})
        return

    slide["layout"] = "comparativo"
    for secao in secoes_genericas:
        itens = secao.get("itens") or [secao["titulo"]]
        slide["listas"].append({"titulo": secao["titulo"], "itens": itens})


def _parse_slide_block(block, idx):
    linhas = [linha.rstrip() for linha in block.splitlines()]
    linhas = [linha for linha in linhas if linha.strip()]
    if not linhas:
        return None

    primeira_linha = linhas[0].strip()
    if primeira_linha.startswith("# "):
        titulo = primeira_linha[2:].strip()
    elif primeira_linha.startswith("## "):
        titulo = primeira_linha[3:].strip()
    else:
        raise ValueError(f"Slide {idx}: o titulo deve iniciar com '# ' ou '## '.")

    titulo_limpo = _remover_prefixo_pagina(_limpar_markdown_inline(titulo))
    if _should_ignore_section(titulo_limpo):
        return None

    slide = {
        "title": titulo_limpo,
        "subtitle": "",
        "layout": "",
        "listas": [],
        "conceitos": [],
    }

    h2_indices = []
    for pos, linha in enumerate(linhas[1:], start=1):
        txt = linha.strip()
        if txt.startswith("## ") and not re.match(r"^##\s+(lista|conceito):", txt, re.IGNORECASE):
            h2_indices.append(pos)

    subtitulo_h2_idx = h2_indices[0] if len(h2_indices) == 1 else None
    secoes_genericas = []
    secao_atual = None
    ignorando_secao = False
    campo_pendente = None

    i = 1
    dentro_codigo = False
    while i < len(linhas):
        linha = linhas[i].strip()

        if linha.startswith("```"):
            dentro_codigo = not dentro_codigo
            i += 1
            continue

        if linha.startswith("> "):
            slide["subtitle"] = _limpar_markdown_inline(linha[2:].strip())
            campo_pendente = None
            i += 1
            continue

        if linha.lower().startswith("layout:"):
            slide["layout"] = linha.split(":", 1)[1].strip().lower()
            campo_pendente = None
            i += 1
            continue

        if subtitulo_h2_idx == i:
            slide["subtitle"] = _limpar_markdown_inline(linha[3:].strip())
            campo_pendente = None
            i += 1
            continue

        m_lista = re.match(r"^##\s+lista:\s*(.+)$", linha, re.IGNORECASE)
        if m_lista:
            titulo_lista = _limpar_markdown_inline(m_lista.group(1).strip())
            itens = []
            i += 1
            while i < len(linhas):
                atual = linhas[i].strip()
                if atual.startswith("## "):
                    break
                if atual.startswith(("- ", "* ")):
                    itens.append(_limpar_markdown_inline(atual[2:].strip()))
                i += 1
            if not itens:
                raise ValueError(f"Slide {idx}, lista '{titulo_lista}': lista sem itens.")
            slide["listas"].append({"titulo": titulo_lista, "itens": itens})
            campo_pendente = None
            continue

        m_conceito = re.match(r"^##\s+conceito:\s*(.+)$", linha, re.IGNORECASE)
        if m_conceito:
            titulo_conceito = _limpar_markdown_inline(m_conceito.group(1).strip())
            descricao_linhas = []
            i += 1
            while i < len(linhas):
                atual = linhas[i].strip()
                if atual.startswith("## "):
                    break
                descricao_linhas.append(_limpar_markdown_inline(atual))
                i += 1
            descricao = " ".join(descricao_linhas).strip()
            if not descricao:
                raise ValueError(f"Slide {idx}, conceito '{titulo_conceito}': descricao vazia.")
            slide["conceitos"].append({"titulo": titulo_conceito, "descricao": descricao})
            campo_pendente = None
            continue

        m_secao = re.match(r"^(##|###|####)\s+(.+)$", linha)
        if m_secao:
            titulo_secao = _remover_prefixo_pagina(_limpar_markdown_inline(m_secao.group(2).strip()))
            if _should_ignore_section(titulo_secao):
                ignorando_secao = True
                secao_atual = None
                i += 1
                continue
            ignorando_secao = False
            rotulo = _normalizar_rotulo(titulo_secao)
            if rotulo in {"titulo", "subtitulo", "descricao"}:
                campo_pendente = rotulo
                secao_atual = None
                i += 1
                continue
            if ":" in titulo_secao:
                prefixo, resto = titulo_secao.split(":", 1)
                if prefixo.strip().lower() in {
                    "coluna esquerda",
                    "coluna direita",
                    "titulo",
                    "subtitulo",
                    "subtítulo",
                    "descrição",
                    "descricao",
                }:
                    titulo_secao = resto.strip()
            secao_atual = {"titulo": titulo_secao, "itens": []}
            secoes_genericas.append(secao_atual)
            campo_pendente = None
            i += 1
            continue

        if ignorando_secao:
            i += 1
            continue

        if linha.startswith(("- ", "* ")):
            texto_item = _limpar_markdown_inline(linha[2:].strip())
            if campo_pendente == "titulo":
                slide["title"] = _remover_prefixo_pagina(texto_item)
                campo_pendente = None
                i += 1
                continue
            if campo_pendente == "subtitulo":
                slide["subtitle"] = texto_item
                campo_pendente = None
                i += 1
                continue
            if not secao_atual:
                secao_atual = {"titulo": "Conteudo", "itens": []}
                secoes_genericas.append(secao_atual)
            m_titulo = re.match(r"^t[íi]tulo:\s*(.+)$", texto_item, re.IGNORECASE)
            if m_titulo:
                secao_atual["titulo"] = m_titulo.group(1).strip()
                i += 1
                continue
            if re.match(r"^descri[cç][aã]o:\s*$", texto_item, re.IGNORECASE):
                i += 1
                continue
            if texto_item:
                secao_atual["itens"].append(texto_item)
            i += 1
            continue

        if linha.startswith("|") and "|" in linha[1:]:
            campo_pendente = None
            if not secao_atual:
                secao_atual = {"titulo": "Tabela", "itens": []}
                secoes_genericas.append(secao_atual)
            secao_atual["itens"].append(_limpar_markdown_inline(linha.replace("|", " ")))
            i += 1
            continue

        if not dentro_codigo and linha:
            if campo_pendente == "titulo":
                slide["title"] = _remover_prefixo_pagina(_limpar_markdown_inline(linha))
                campo_pendente = None
                i += 1
                continue
            if campo_pendente == "subtitulo":
                slide["subtitle"] = _limpar_markdown_inline(linha)
                campo_pendente = None
                i += 1
                continue
            campo_pendente = None
            if not secao_atual:
                secao_atual = {"titulo": "Conteudo", "itens": []}
                secoes_genericas.append(secao_atual)
            secao_atual["itens"].append(_limpar_markdown_inline(linha))

        i += 1

    if not slide["layout"]:
        _inferir_layout_generico(slide, secoes_genericas)

    if not slide["layout"]:
        layouts = ", ".join(sorted(LAYOUTS_SUPORTADOS))
        raise ValueError(
            f"Slide {idx}: nao foi possivel inferir o layout. "
            f"Use um layout explicito, por exemplo: {layouts}."
        )
    if slide["layout"] not in LAYOUTS_SUPORTADOS:
        layouts = ", ".join(sorted(LAYOUTS_SUPORTADOS))
        raise ValueError(f"Slide {idx}: layout invalido '{slide['layout']}'. Use: {layouts}.")
    if slide["layout"] == "comparativo" and not slide["listas"]:
        raise ValueError(f"Slide {idx}: layout comparativo exige ao menos uma secao '## lista:'.")
    if slide["layout"] == "conceitos" and not slide["conceitos"]:
        raise ValueError(f"Slide {idx}: layout conceitos exige ao menos uma secao '## conceito:'.")

    return slide


def carregar_slides_markdown(caminho_markdown):
    """
    Le o arquivo Markdown e retorna uma lista de slides.

    Contrato principal:
    - Cada slide comeca com '# Titulo'
    - Slides podem ser separados por '---' ou '<!-- SLIDE -->'
    - Subtitulo opcional com '> ...'
    - Layout com 'layout: comparativo|conceitos' (ou inferencia automatica)
    - Listas com '## lista: Titulo' + itens '- ...'
    - Conceitos com '## conceito: Titulo' + paragrafo de descricao
    """
    conteudo = Path(caminho_markdown).read_text(encoding="utf-8")
    blocos = [
        b.strip()
        for b in re.split(r"^\s*(?:---|<!--\s*SLIDE\s*-->)\s*$", conteudo, flags=re.MULTILINE)
    ]
    slides = []
    for idx, bloco in enumerate(blocos, start=1):
        if not bloco:
            continue
        slide = _parse_slide_block(bloco, idx)
        if slide:
            slides.append(slide)
    if not slides:
        raise ValueError("Nenhum slide valido encontrado no arquivo Markdown.")
    return slides


def render_slide_comparativo(slide, data):
    layout = THEME_ACTIVE["layout"]
    slide_cfg = THEME_ACTIVE["slide"]
    mx = layout["margin_x_in"]
    gap = layout["column_gap_in"]
    content_top = layout["content_top_in"]
    bottom_top = layout["bottom_block_top_in"]
    bottom_height = layout["bottom_block_height_in"]
    usable_width = slide_cfg["width_in"] - (mx * 2)
    col_width = (usable_width - gap) / 2

    posicoes = [
        (Inches(mx), Inches(content_top), Inches(col_width), Inches(2.4), 20, 14),
        (Inches(mx + col_width + gap), Inches(content_top), Inches(col_width), Inches(2.4), 20, 14),
        (Inches(mx), Inches(bottom_top), Inches(usable_width), Inches(bottom_height), 18, 13),
    ]

    for idx, secao in enumerate(data["listas"]):
        if idx < len(posicoes):
            left, top, width, height, ft_titulo, ft_item = posicoes[idx]
        else:
            left = Inches(mx)
            top = Inches(bottom_top + 0.2 + (idx - 2) * 1.2)
            width = Inches(usable_width)
            height = Inches(1.5)
            ft_titulo = 16
            ft_item = 12

        adicionar_lista(
            slide,
            itens=secao["itens"],
            left=left,
            top=top,
            width=width,
            height=height,
            titulo=secao["titulo"],
            fonte_titulo=ft_titulo,
            fonte_item=ft_item,
        )


def render_slide_conceitos(slide, data):
    layout = THEME_ACTIVE["layout"]
    slide_cfg = THEME_ACTIVE["slide"]
    mx = layout["margin_x_in"]
    gap = layout["column_gap_in"]
    content_top = layout["content_top_in"] - 0.1
    bottom_top = layout["bottom_block_top_in"]
    usable_width = slide_cfg["width_in"] - (mx * 2)
    col_width = (usable_width - gap) / 2

    posicoes_conceitos = [
        (Inches(mx), Inches(content_top), Inches(col_width), Inches(1.9)),
        (Inches(mx + col_width + gap), Inches(content_top), Inches(col_width), Inches(1.9)),
    ]

    for idx, conceito in enumerate(data["conceitos"]):
        if idx < len(posicoes_conceitos):
            left, top, width, height = posicoes_conceitos[idx]
        else:
            left = Inches(mx)
            top = Inches(3.9 + (idx - 2) * 1.8)
            width = Inches(usable_width)
            height = Inches(1.6)
        adicionar_caixa_conceito(
            slide,
            titulo=conceito["titulo"],
            descricao=conceito["descricao"],
            left=left,
            top=top,
            width=width,
            height=height,
        )

    if data["listas"]:
        adicionar_lista(
            slide,
            itens=[f"{item}" for secao in data["listas"] for item in secao["itens"]],
            left=Inches(mx),
            top=Inches(bottom_top),
            width=Inches(usable_width),
            height=Inches(2.2),
            titulo="Anatomia do Agente",
            fonte_titulo=18,
            fonte_item=12,
        )


def _add_body_lines(slide, lines, left, top, width, height, size=None):
    size = size or THEME_ACTIVE["fonts"]["body_size"]
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    _ajustar_text_frame(tf)
    first = True
    for line in lines:
        if not line:
            continue
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = line
        run.font.name = FONTE_PRINCIPAL
        run.font.size = Pt(size)
        run.font.color.rgb = COR_TEXTO
    return box


def render_slide_title_left_text_right(slide, data):
    layout = THEME_ACTIVE["layout"]
    slide_cfg = THEME_ACTIVE["slide"]
    mx = layout["margin_x_in"]
    top = layout.get("max_content_top_in", 2.2)
    left_width = layout.get("max_left_title_width_in", 3.2)
    usable_width = slide_cfg["width_in"] - (mx * 2)
    gap = layout["column_gap_in"]
    right_x = mx + left_width + gap
    right_width = max(usable_width - left_width - gap, 3.0)
    lines = []
    for secao in data["listas"]:
        if secao.get("titulo"):
            lines.append(secao["titulo"])
        lines.extend(secao.get("itens", []))
        lines.append("")
    _add_body_lines(
        slide,
        lines[:12],
        Inches(right_x),
        Inches(top),
        Inches(right_width),
        Inches(4.7),
        size=max(THEME_ACTIVE["fonts"]["body_size"], 18),
    )


def render_slide_title_top_bullets(slide, data):
    mx = THEME_ACTIVE["layout"]["margin_x_in"]
    top = THEME_ACTIVE["layout"].get("max_content_top_in", 2.2)
    width = THEME_ACTIVE["slide"]["width_in"] - (mx * 2)
    y = top
    max_blocos = 3
    for secao in data["listas"][:max_blocos]:
        titulo = secao.get("titulo", "")
        itens = secao.get("itens", [])
        clean_items = []
        titulo_norm = _limpar_markdown_inline(titulo).strip().lower().rstrip(":")
        for item in itens:
            val = _limpar_markdown_inline(item).strip()
            if not val:
                continue
            val_norm = val.lower().rstrip(":")
            if val_norm == titulo_norm:
                continue
            if val_norm.startswith(f"{titulo_norm}:"):
                continue
            clean_items.append(val)

        lines = [f"• {item}" for item in clean_items[:3]]
        box = slide.shapes.add_textbox(Inches(mx), Inches(y), Inches(width), Inches(1.55))
        tf = box.text_frame
        tf.clear()
        _ajustar_text_frame(tf)
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = titulo
        run.font.name = FONTE_TITULO
        run.font.size = Pt(20)
        run.font.color.rgb = COR_TITULO
        run.font.bold = False

        for line in lines:
            p2 = tf.add_paragraph()
            r2 = p2.add_run()
            r2.text = line
            r2.font.name = FONTE_PRINCIPAL
            r2.font.size = Pt(15)
            r2.font.color.rgb = COR_TEXTO

        y += 1.45


def render_slide_title_top_grid_2x2(slide, data):
    mx = THEME_ACTIVE["layout"]["margin_x_in"]
    top = THEME_ACTIVE["layout"].get("max_content_top_in", 2.2)
    width = THEME_ACTIVE["slide"]["width_in"] - (mx * 2)
    gap = THEME_ACTIVE["layout"]["column_gap_in"]
    col_width = (width - gap) / 2
    row_height = 2.05
    cards = data["listas"][:4]
    for idx, secao in enumerate(cards):
        col = idx % 2
        row = idx // 2
        left = mx + col * (col_width + gap)
        top_pos = top + row * row_height
        header = f"{idx + 1:02d}"
        body_title = secao.get("titulo", "")
        body_text = " ".join(secao.get("itens", [])[:2])
        _add_body_lines(slide, [header], Inches(left), Inches(top_pos), Inches(col_width), Inches(0.35), size=16)
        _add_body_lines(slide, [body_title], Inches(left), Inches(top_pos + 0.35), Inches(col_width), Inches(0.55), size=18)
        _add_body_lines(slide, [body_text], Inches(left), Inches(top_pos + 0.9), Inches(col_width), Inches(0.95), size=15)


def render_slide_title_top_text_block(slide, data):
    mx = THEME_ACTIVE["layout"]["margin_x_in"]
    top = THEME_ACTIVE["layout"].get("max_content_top_in", 2.2)
    width = THEME_ACTIVE["slide"]["width_in"] - (mx * 2)
    text_parts = []
    for secao in data["listas"]:
        if secao.get("itens"):
            text_parts.append(" ".join(secao["itens"]))
        elif secao.get("titulo"):
            text_parts.append(secao["titulo"])
    _add_body_lines(
        slide,
        [" ".join(text_parts)],
        Inches(mx + (width * 0.2)),
        Inches(top),
        Inches(width * 0.76),
        Inches(3.7),
        size=15,
    )


def renderizar_apresentacao(slides_data):
    prs = Presentation()
    prs.slide_width = Inches(THEME_ACTIVE["slide"]["width_in"])
    prs.slide_height = Inches(THEME_ACTIVE["slide"]["height_in"])
    for data in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        configurar_slide(slide)
        adicionar_titulo(slide, data["title"])
        if data["subtitle"]:
            adicionar_subtitulo(slide, data["subtitle"])

        if data["layout"] == "comparativo":
            render_slide_comparativo(slide, data)
        elif data["layout"] == "conceitos":
            render_slide_conceitos(slide, data)
        elif data["layout"] == "titulo":
            pass
        elif data["layout"] == "title_left_text_right":
            render_slide_title_left_text_right(slide, data)
        elif data["layout"] == "title_top_bullets":
            render_slide_title_top_bullets(slide, data)
        elif data["layout"] == "title_top_grid_2x2":
            render_slide_title_top_grid_2x2(slide, data)
        elif data["layout"] == "title_top_text_block":
            render_slide_title_top_text_block(slide, data)
        else:
            print(f"Aviso: layout desconhecido '{data['layout']}'. Slide ignorado.")

    return prs


def parse_args():
    parser = argparse.ArgumentParser(description="Gerador de slides PowerPoint a partir de Markdown.")
    parser.add_argument(
        "--input",
        default="slides.md",
        help="Arquivo markdown de entrada (padrao: slides.md).",
    )
    parser.add_argument(
        "--output",
        default="Slides_Agentic_AI.pptx",
        help="Arquivo .pptx de saida (padrao: Slides_Agentic_AI.pptx).",
    )
    parser.add_argument(
        "--theme",
        default=None,
        help="Arquivo JSON de tema visual (opcional).",
    )
    parser.add_argument(
        "--profile",
        default=None,
        choices=["default", "premium", "max"],
        help="Perfil visual predefinido (opcional).",
    )
    return parser.parse_args()


def main():
    args = parse_args()
    try:
        if args.theme and not Path(args.theme).exists():
            raise ValueError(f"arquivo de tema nao encontrado: '{args.theme}'")
        if not Path(args.input).exists():
            raise FileNotFoundError(args.input)
        tema = _load_theme(args.theme, args.profile)
        _apply_theme(tema)
        slides_data = carregar_slides_markdown(args.input)
        apresentacao = renderizar_apresentacao(slides_data)
        apresentacao.save(args.output)
        print(f"Slides criados com sucesso: '{args.output}'")
    except FileNotFoundError:
        print(f"Erro: arquivo de entrada nao encontrado: '{args.input}'")
        print("Dica: verifique o nome do arquivo ou passe o caminho completo com --input.")
        sys.exit(1)
    except ValueError as exc:
        print(f"Erro de validacao do markdown: {exc}")
        sys.exit(1)


if __name__ == "__main__":
    main()